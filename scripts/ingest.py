#!/usr/bin/env python3
"""Ingesta documental para el repositorio Micromódulo IA Salud.

Recorre ``raw_data/`` (o un directorio alternativo), extrae texto y metadatos
de PDF / DOCX / PPTX / XLSX / MD / TXT, genera markdown normalizado en
``content/`` con frontmatter YAML, y actualiza ``data/document-registry.yaml``.

Diseño:

* **Idempotente**: si el hash SHA-256 del fichero fuente coincide con una
  entrada del registro, no se vuelve a generar el markdown salvo
  ``--force``.
* **No destructivo**: nunca borra ni sobrescribe páginas markdown previas
  con contenido distinto; en su lugar crea un sufijo ``-conflict`` y avisa.
* **Detección de familia**: agrupa ficheros por *stem* normalizado
  (quitando sufijos de versión obvios como ``_v2``, ``_REV``,
  ``_FINAL``, fechas ``2024``, ``rev JV``, etc.). Es heurística;
  un humano lo confirma.

Uso típico:

    python scripts/ingest.py
    python scripts/ingest.py --raw raw_data/master-training --force
    python scripts/ingest.py --dry-run

Salida:

* ``content/<slug>.md`` por documento.
* ``data/document-registry.yaml`` actualizado.
* ``ingestion/run-<timestamp>.log`` con el log de la pasada.
"""
from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import logging
import os
import re
import sys
import unicodedata
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import Iterable

import yaml

# Optional/heavy deps imported lazily inside extractors so a single missing
# dependency doesn't break the whole script.

# --------------------------------------------------------------------------- #
# Paths
# --------------------------------------------------------------------------- #

# Resolved at runtime; the script lives in <repo>/documentation/scripts/.
SCRIPT_DIR = Path(__file__).resolve().parent
DOC_ROOT = SCRIPT_DIR.parent                       # documentation/
PROFILE_ROOT = DOC_ROOT.parent                     # master-healthcare-ai profile
DEFAULT_RAW = PROFILE_ROOT / "raw_data"
CONTENT_DIR = DOC_ROOT / "content"
REGISTRY_PATH = DOC_ROOT / "data" / "document-registry.yaml"
INGEST_LOG_DIR = DOC_ROOT / "ingestion"

CONTENT_DIR.mkdir(parents=True, exist_ok=True)
REGISTRY_PATH.parent.mkdir(parents=True, exist_ok=True)
INGEST_LOG_DIR.mkdir(parents=True, exist_ok=True)

# --------------------------------------------------------------------------- #
# Logging
# --------------------------------------------------------------------------- #

def _setup_logging() -> Path:
    ts = dt.datetime.now().strftime("%Y%m%d-%H%M%S")
    log_path = INGEST_LOG_DIR / f"run-{ts}.log"
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(levelname)s %(message)s",
        handlers=[
            logging.FileHandler(log_path, encoding="utf-8"),
            logging.StreamHandler(sys.stdout),
        ],
    )
    return log_path


# --------------------------------------------------------------------------- #
# Utilities
# --------------------------------------------------------------------------- #

# Filenames we never ingest.
SKIP_NAMES = {".DS_Store", "Thumbs.db", "desktop.ini"}
SKIP_SUFFIXES = {".tmp", ".lock"}


def slugify(value: str, maxlen: int = 80) -> str:
    """Filename-safe slug. ASCII, lower, hyphen-separated."""
    value = unicodedata.normalize("NFKD", value)
    value = value.encode("ascii", "ignore").decode("ascii")
    value = re.sub(r"[^a-zA-Z0-9]+", "-", value).strip("-").lower()
    return value[:maxlen] or "doc"


# Patterns we strip when computing a "family stem" — the goal is to recognise
# that "Foo v1.docx", "Foo_v2.docx", "Foo - rev JV.docx", "Foo (final).pdf"
# all belong to the same family.
_FAMILY_STRIP_PATTERNS = [
    r"\bv\d+(\.\d+)?\b",
    r"\brev(?:\s|_)?\w*\b",
    r"\b(?:final|definitivo|aprobado|approved|signed|firmado)\d*\b",
    r"\bcopy\s*of\b",
    r"\bdraft\b",
    r"\b20\d{2}(?:[-_]?\d{2})?(?:[-_]?\d{2})?\b",   # dates / years
    r"\b\d{8}\b",                                    # YYYYMMDD
    r"\b(?:dic|ene|feb|mar|abr|may|jun|jul|ago|sep|oct|nov)\s*\d*\b",
    r"\s+jv\b",
    r"\s+jennifer\b",
    r"\b_+\b",
]


def family_stem(name: str) -> str:
    """Heuristic stem used to group versions of the same document.

    Tested against the corpus in ``raw_data/master-training/``; refine as we
    discover edge cases.
    """
    stem = Path(name).stem
    stem = stem.lower()
    stem = unicodedata.normalize("NFKD", stem).encode("ascii", "ignore").decode("ascii")
    for pat in _FAMILY_STRIP_PATTERNS:
        stem = re.sub(pat, " ", stem, flags=re.IGNORECASE)
    stem = re.sub(r"[\W_]+", " ", stem).strip()
    stem = re.sub(r"\s+", " ", stem)
    return stem


def file_sha256(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def safe_yaml_dump(obj) -> str:
    return yaml.safe_dump(obj, sort_keys=False, allow_unicode=True, width=100)


# --------------------------------------------------------------------------- #
# Extractors
# --------------------------------------------------------------------------- #

@dataclass
class ExtractResult:
    text: str
    extra_meta: dict = field(default_factory=dict)
    n_pages: int | None = None


def extract_pdf(path: Path) -> ExtractResult:
    import fitz  # pymupdf
    doc = fitz.open(path)
    chunks: list[str] = []
    for i, page in enumerate(doc, 1):
        chunks.append(f"\n\n<!-- página {i} -->\n\n")
        chunks.append(page.get_text("text"))
    meta = {k: v for k, v in (doc.metadata or {}).items() if v}
    return ExtractResult(text="".join(chunks), extra_meta={"pdf_meta": meta}, n_pages=len(doc))


def extract_docx(path: Path) -> ExtractResult:
    import docx
    d = docx.Document(path)
    lines: list[str] = []
    for p in d.paragraphs:
        style = (p.style.name or "").lower() if p.style else ""
        txt = p.text.strip()
        if not txt:
            lines.append("")
            continue
        if style.startswith("heading 1"):
            lines.append(f"# {txt}")
        elif style.startswith("heading 2"):
            lines.append(f"## {txt}")
        elif style.startswith("heading 3"):
            lines.append(f"### {txt}")
        elif style.startswith("heading"):
            lines.append(f"#### {txt}")
        else:
            lines.append(txt)
    # tables
    for ti, table in enumerate(d.tables, 1):
        lines.append("")
        lines.append(f"<!-- tabla {ti} -->")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
        lines.append("")
    return ExtractResult(text="\n".join(lines))


def extract_pptx(path: Path) -> ExtractResult:
    from pptx import Presentation
    prs = Presentation(path)
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n## Diapositiva {i}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append(f"\n> **Notas:** {notes}\n")
    return ExtractResult(text="\n".join(lines), n_pages=len(prs.slides))


def extract_xlsx(path: Path) -> ExtractResult:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    lines: list[str] = []
    for sheet in wb.worksheets:
        lines.append(f"\n## Hoja: {sheet.title}\n")
        rows = list(sheet.iter_rows(values_only=True))
        # Trim trailing empty rows / cols
        rows = [list(r) for r in rows if any(c is not None and str(c).strip() for c in r)]
        if not rows:
            lines.append("_(hoja vacía)_")
            continue
        # Render as markdown table — but cap width to avoid 200-col monsters
        max_cols = min(max(len(r) for r in rows), 20)
        header = rows[0][:max_cols]
        lines.append("| " + " | ".join(str(c) if c is not None else "" for c in header) + " |")
        lines.append("|" + "|".join(["---"] * max_cols) + "|")
        for r in rows[1:200]:        # cap rows in preview
            r = r[:max_cols]
            lines.append("| " + " | ".join(str(c) if c is not None else "" for c in r) + " |")
        if len(rows) > 200:
            lines.append(f"\n_... ({len(rows) - 200} filas adicionales truncadas en la vista previa)_")
    return ExtractResult(text="\n".join(lines))


def extract_text(path: Path) -> ExtractResult:
    return ExtractResult(text=path.read_text(encoding="utf-8", errors="replace"))


EXTRACTORS: dict[str, callable] = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".md": extract_text,
    ".txt": extract_text,
}


# --------------------------------------------------------------------------- #
# Registry
# --------------------------------------------------------------------------- #

@dataclass
class RegistryEntry:
    id: str
    title: str
    category: str
    version: str
    status: str
    author: str
    created: str
    updated: str
    supersedes: list = field(default_factory=list)
    superseded_by: str = ""
    source_file: str = ""
    source_sha256: str = ""
    family_id: str = ""
    canonical: bool = False
    tags: list = field(default_factory=list)
    content_path: str = ""
    language: str = "es"
    n_pages: int | None = None


def load_registry() -> dict:
    if REGISTRY_PATH.exists():
        with REGISTRY_PATH.open("r", encoding="utf-8") as f:
            data = yaml.safe_load(f) or {}
    else:
        data = {}
    data.setdefault("documents", [])
    data.setdefault("families", {})
    data.setdefault("meta", {})
    return data


def save_registry(registry: dict) -> None:
    registry["meta"]["last_run"] = dt.datetime.now(dt.timezone.utc).isoformat()
    registry["meta"]["doc_count"] = len(registry["documents"])
    with REGISTRY_PATH.open("w", encoding="utf-8") as f:
        f.write("# Registro canónico de documentos. Generado por scripts/ingest.py.\n")
        f.write("# La revisión humana decide canonical/status/family_id finales.\n\n")
        yaml.safe_dump(registry, f, sort_keys=False, allow_unicode=True, width=100)


def next_doc_id(registry: dict) -> str:
    n = 1
    used = {entry["id"] for entry in registry["documents"] if entry["id"].startswith("DOC-")}
    while f"DOC-{n:04d}" in used:
        n += 1
    return f"DOC-{n:04d}"


# --------------------------------------------------------------------------- #
# Category inference
# --------------------------------------------------------------------------- #

def infer_category(rel_path: str, name: str) -> str:
    p = rel_path.lower()
    n = name.lower()
    if "micromodulo" in p or "micromódulo" in p or "micromodulo" in n or "micromódulo" in n:
        return "micromodulo-ia-salud"
    if "telemedicina" in p:
        return "asignatura-telemedicina"
    if "especialista" in p or "informática de la salud" in p or "informatica de la salud" in p:
        return "curso-especialista-informatica-salud"
    if "convenio" in n:
        return "administracion-convenios"
    if "factura" in n or "financiero" in p:
        return "administracion-financiero"
    if "contrato" in n:
        return "administracion-contratos"
    if "anexo" in n or "propuesta" in n:
        return "propuestas-academicas"
    if "actas" in p or "acta" in n:
        return "actas"
    if "cv" in n or "ficha profesor" in n:
        return "faculty"
    if "marketing" in n or "folleto" in n:
        return "marketing"
    return "otros"


def infer_tags(rel_path: str, name: str) -> list[str]:
    tags = []
    p = (rel_path + " " + name).lower()
    for kw, tag in [
        ("micromódulo", "micromodulo"),
        ("micromodulo", "micromodulo"),
        ("inteligencia artificial", "ia"),
        ("telemedicina", "telemedicina"),
        ("especialista", "especialista"),
        ("convenio", "convenio"),
        ("factura", "financiero"),
        ("contrato", "contrato"),
        ("anexo i", "anexo-i"),
        ("anexo ii", "anexo-ii"),
        ("anexo iii", "anexo-iii"),
        ("marketing", "marketing"),
        ("folleto", "marketing"),
        ("cv", "cv"),
        ("salud", "salud"),
        ("edge", "edgehealth"),
    ]:
        if kw in p and tag not in tags:
            tags.append(tag)
    return tags


# --------------------------------------------------------------------------- #
# Pipeline
# --------------------------------------------------------------------------- #

def iter_source_files(raw_root: Path) -> Iterable[Path]:
    raw_root = raw_root.resolve()
    for p in sorted(raw_root.rglob("*")):
        if not p.is_file():
            continue
        p = p.resolve()
        if p.name in SKIP_NAMES or p.suffix.lower() in SKIP_SUFFIXES:
            continue
        if p.suffix.lower() not in EXTRACTORS and p.suffix.lower() not in {".doc"}:
            # Unknown extension — we still want to register it, just without text.
            yield p
            continue
        yield p


def render_markdown(entry: RegistryEntry, extract: ExtractResult | None, source_rel: str) -> str:
    fm = {
        "id": entry.id,
        "title": entry.title,
        "version": entry.version,
        "status": entry.status,
        "authors": [entry.author] if entry.author else [],
        "created": entry.created,
        "updated": entry.updated,
        "tags": entry.tags,
        "related": [],
        "source": source_rel,
        "source_sha256": entry.source_sha256,
        "canonical": entry.canonical,
        "family_id": entry.family_id,
        "supersedes": entry.supersedes,
        "superseded_by": entry.superseded_by,
        "language": entry.language,
        "category": entry.category,
        "n_pages": entry.n_pages,
    }
    body = []
    body.append("---")
    body.append(safe_yaml_dump(fm).rstrip())
    body.append("---")
    body.append("")
    body.append(f"# {entry.title}")
    body.append("")
    body.append(f"> Generado automáticamente desde `{source_rel}` por `scripts/ingest.py`. "
                f"Revisar metadatos en el frontmatter y en `data/document-registry.yaml`.")
    body.append("")
    if extract is None:
        body.append("_Formato no extraído automáticamente (binario sin extractor)._")
        body.append("")
    else:
        text = extract.text.strip()
        # Hard cap to keep markdown reasonable; the source is always the truth.
        if len(text) > 200_000:
            text = text[:200_000] + "\n\n_... (truncado a 200 000 caracteres; ver fichero fuente)_"
        body.append(text)
    return "\n".join(body) + "\n"


def process_file(
    src: Path,
    raw_root: Path,
    registry: dict,
    force: bool,
    dry_run: bool,
) -> RegistryEntry | None:
    try:
        rel = src.resolve().relative_to(PROFILE_ROOT).as_posix()
    except ValueError:
        # Source outside the profile root — store the absolute path instead.
        rel = src.resolve().as_posix()
    sha = file_sha256(src)

    # Idempotency: skip if same source already present (unless --force).
    existing = next(
        (e for e in registry["documents"] if e.get("source_sha256") == sha),
        None,
    )
    if existing and not force:
        logging.info("skip (already ingested) %s", rel)
        return None

    suffix = src.suffix.lower()
    extractor = EXTRACTORS.get(suffix)

    extract = None
    n_pages = None
    if extractor:
        try:
            extract = extractor(src)
            n_pages = extract.n_pages
        except Exception as e:
            logging.warning("extraction failed for %s: %s", rel, e)
            extract = ExtractResult(text=f"_Error de extracción: {e}_")
    elif suffix == ".doc":
        # Legacy .doc — no native extractor. Register but mark for manual handling.
        logging.warning("legacy .doc (sin extractor): %s", rel)
    else:
        logging.info("unknown ext %s — registering metadata only", rel)

    # Build identity
    if existing:
        entry = RegistryEntry(**{k: existing.get(k, v) for k, v in asdict(_blank_entry()).items()})
        entry.id = existing["id"]
        entry.updated = dt.date.today().isoformat()
    else:
        entry = _blank_entry()
        entry.id = next_doc_id(registry)
        # Reserve the ID immediately so concurrent files get distinct IDs.
        registry["documents"].append({"id": entry.id, "_reserved": True})

    title = src.stem.strip()
    entry.title = title
    entry.category = infer_category(rel, src.name)
    entry.tags = infer_tags(rel, src.name)
    entry.source_file = rel
    entry.source_sha256 = sha
    entry.family_id = "fam-" + slugify(family_stem(src.name))
    entry.n_pages = n_pages
    if not entry.created:
        try:
            entry.created = dt.date.fromtimestamp(src.stat().st_mtime).isoformat()
        except Exception:
            entry.created = dt.date.today().isoformat()
    entry.updated = dt.date.today().isoformat()

    # Markdown destination
    slug = slugify(f"{entry.id}-{title}")
    dest = CONTENT_DIR / f"{slug}.md"
    entry.content_path = dest.relative_to(DOC_ROOT).as_posix()

    md = render_markdown(entry, extract, rel)

    if dry_run:
        logging.info("DRY-RUN would write %s", dest)
    else:
        if dest.exists():
            existing_md = dest.read_text(encoding="utf-8", errors="replace")
            if existing_md != md and not force:
                # Conflict: write to a -conflict path and warn.
                conflict = dest.with_name(dest.stem + "-conflict.md")
                conflict.write_text(md, encoding="utf-8")
                logging.warning("conflict — wrote %s instead of overwriting %s",
                                conflict, dest)
            else:
                dest.write_text(md, encoding="utf-8")
        else:
            dest.write_text(md, encoding="utf-8")

    # Persist back to registry (replace the reservation or update existing)
    persisted = asdict(entry)
    docs = registry["documents"]
    # Remove the reservation we inserted (or any existing entry with same id)
    docs[:] = [d for d in docs if d.get("id") != entry.id]
    docs.append(persisted)

    # Family bookkeeping
    fam = registry["families"].setdefault(entry.family_id, {"members": [], "canonical": None})
    if entry.id not in fam["members"]:
        fam["members"].append(entry.id)

    logging.info("ingested %s  →  %s  (family: %s)", rel, entry.content_path, entry.family_id)
    return entry


def _blank_entry() -> RegistryEntry:
    today = dt.date.today().isoformat()
    return RegistryEntry(
        id="",
        title="",
        category="otros",
        version="1.0",
        status="Draft",
        author="",
        created=today,
        updated=today,
    )


# --------------------------------------------------------------------------- #
# Main
# --------------------------------------------------------------------------- #

def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    parser.add_argument("--raw", type=Path, default=DEFAULT_RAW,
                        help=f"Raíz de datos crudos (por defecto {DEFAULT_RAW}).")
    parser.add_argument("--force", action="store_true",
                        help="Regenera markdown aunque el hash coincida.")
    parser.add_argument("--dry-run", action="store_true",
                        help="No escribe ficheros; sólo loguea lo que haría.")
    args = parser.parse_args()

    log_path = _setup_logging()
    logging.info("ingest start  raw=%s  force=%s  dry_run=%s",
                 args.raw, args.force, args.dry_run)
    logging.info("log file: %s", log_path)

    if not args.raw.exists():
        logging.error("raw root no existe: %s", args.raw)
        return 2

    registry = load_registry()
    n_seen = n_ingested = n_skipped = 0

    for src in iter_source_files(args.raw):
        n_seen += 1
        try:
            result = process_file(src, args.raw, registry, args.force, args.dry_run)
        except Exception as e:
            logging.exception("error processing %s: %s", src, e)
            continue
        if result is None:
            n_skipped += 1
        else:
            n_ingested += 1

    # Drop any leftover reservations (defensive)
    registry["documents"] = [d for d in registry["documents"] if not d.get("_reserved")]

    if not args.dry_run:
        save_registry(registry)

    logging.info("summary  seen=%d  ingested=%d  skipped=%d  total_in_registry=%d",
                 n_seen, n_ingested, n_skipped, len(registry["documents"]))
    return 0


if __name__ == "__main__":
    sys.exit(main())
